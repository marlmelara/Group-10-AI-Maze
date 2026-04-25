"""Generate presentation.pptx from the LaTeX deck's content.

PowerPoint natively plays animated GIFs when a GIF is embedded as a
picture, so this deck can actually animate the beta_demo GIFs during
slideshow mode.

Usage: python3 src/make_pptx.py
Output: presentation.pptx
"""
from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "figures"
OUT = ROOT / "presentation.pptx"

# 16:9 standard widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BRAND_NAVY = RGBColor(0x1F, 0x3A, 0x5F)
BRAND_LIGHT = RGBColor(0xF6, 0xF5, 0xEE)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)


def add_title_block(slide, text: str, subtitle: str = ""):
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND_NAVY
    bar.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.4), Inches(0.12),
                                    SLIDE_W - Inches(0.8), Inches(0.7))
    tf = tbox.text_frame
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.4), Inches(0.95),
                                      SLIDE_W - Inches(0.8), Inches(0.4))
        p2 = sb.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.italic = True
        p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)


def add_bullets(slide, items, left: float, top: float, width: float,
                height: float, font_size: int = 18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + text
        p.font.size = Pt(font_size)
        p.space_after = Pt(6)


def add_footer(slide, page_num: int, total: int):
    box = slide.shapes.add_textbox(Inches(0.4), SLIDE_H - Inches(0.35),
                                   SLIDE_W - Inches(0.8), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = (f"Group 15  •  Marlon Melara, Ahnaf Murshid  |  "
              f"Silent Cartographer: Hybrid RL + A* Agent  |  "
              f"April 22, 2026  |  {page_num}/{total}")
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p.alignment = PP_ALIGN.CENTER


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    BLANK = prs.slide_layouts[6]

    total = 13

    # ------------------- 1. Title -------------------
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRAND_NAVY
    bg.line.fill.background()
    t = s.shapes.add_textbox(Inches(0.5), Inches(2.2),
                             SLIDE_W - Inches(1.0), Inches(1.6))
    p = t.text_frame.paragraphs[0]
    p.text = "Silent Cartographer"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    t = s.shapes.add_textbox(Inches(0.5), Inches(3.6),
                             SLIDE_W - Inches(1.0), Inches(0.7))
    p = t.text_frame.paragraphs[0]
    p.text = "Hybrid Reinforcement Learning + A* Maze Agent"
    p.font.size = Pt(26)
    p.font.color.rgb = RGBColor(0xE6, 0xE6, 0xE6)
    p.alignment = PP_ALIGN.CENTER

    t = s.shapes.add_textbox(Inches(0.5), Inches(5.2),
                             SLIDE_W - Inches(1.0), Inches(1.2))
    tf = t.text_frame
    for i, line in enumerate([
        "COSC 4368  •  Spring 2026  •  Final Submission",
        "Group 15 — Marlon Melara, Ahnaf Murshid",
        "April 22, 2026",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0xD8, 0xD8, 0xD8)
        p.alignment = PP_ALIGN.CENTER

    # ------------------- 2. Task Recap -------------------
    s = prs.slides.add_slide(BLANK)
    add_title_block(s, "Task Recap")
    add_bullets(s, [
        "64×64 grid maze, agent has no sensors — learns only by bumping and dying.",
        "Hazards: rotating V-shaped fire clusters (90° every 5 actions), teleport pads, confusion pads.",
        "Each turn, submit 1–5 actions; environment returns a TurnResult.",
        "Train on maze-alpha; test on maze-beta (no retraining per spec).",
        "Required metrics: success rate, avg turns, avg path length, death rate.",
        "Bonus metrics: exploration efficiency, map completeness, replanning efficiency, learning efficiency.",
    ], left=0.6, top=1.4, width=12.1, height=5.5, font_size=20)
    add_footer(s, 2, total)

    # ------------------- 3. Approach -------------------
    s = prs.slides.add_slide(BLANK)
    add_title_block(s, "Approach — Why Hybrid RL + A*")
    # Left column
    tbox = s.shapes.add_textbox(Inches(0.6), Inches(1.4),
                                Inches(6.0), Inches(5.5))
    tf = tbox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Reinforcement Learning"
    p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = ACCENT
    for text in ["Handles delayed, sparse hazard feedback — you only learn a cell is a fire after you die on it.",
                 "TD(0) updates spread danger to Manhattan neighbors, so the agent detours regions, not cells.",
                 "Hyperparameters transfer alpha → beta even when cell state resets."]:
        p = tf.add_paragraph()
        p.text = "• " + text
        p.font.size = Pt(17); p.space_after = Pt(6)
    # Right column
    tbox = s.shapes.add_textbox(Inches(7.0), Inches(1.4),
                                Inches(5.7), Inches(5.5))
    tf = tbox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Classical A* Search"
    p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = ACCENT
    for text in ["Known edges are deterministic forever — search exploits that immediately.",
                 "Manhattan heuristic keeps 64×64 planning in milliseconds.",
                 "Visit counts + risk penalties become edge weights in phase-aware BFS."]:
        p = tf.add_paragraph()
        p.text = "• " + text
        p.font.size = Pt(17); p.space_after = Pt(6)
    add_footer(s, 3, total)

    # ------------------- 4. Agent Architecture -------------------
    s = prs.slides.add_slide(BLANK)
    add_title_block(s, "Agent Architecture")
    add_bullets(s, [
        "Loop: environment returns TurnResult → agent integrates → plans → submits 1–5 actions.",
        "Integration updates: walls, per-phase fire sets, teleport map, confusion pads, Q-table.",
        "Planner: phase-aware BFS + goal-biased frontier selector.",
        "RL layer: 64×64×4 Q-table decayed between episodes (prevents stale fear).",
        "Same loop, same planner, same hyperparameters across all three mazes.",
    ], left=0.6, top=1.4, width=12.1, height=5.5, font_size=20)
    add_footer(s, 4, total)

    # ------------------- 5. Rotating Fire Hazards -------------------
    s = prs.slides.add_slide(BLANK)
    add_title_block(s, "Rotating Fire Hazards")
    # Text left
    add_bullets(s, [
        "Fires arranged in V-clusters, pivot at the tip of the V.",
        "Rotate 90° clockwise every 5 cumulative actions.",
        "Agent has no visual — learns by dying.",
        "We pre-compute all 4 rotation states at environment construction.",
        "Conservative: a cell that has EVER been a fire is dangerous across all phases by default.",
    ], left=0.6, top=1.4, width=6.5, height=5.5, font_size=18)
    # Image right
    parse_img = FIG / "parse_check_beta.png"
    if parse_img.exists():
        s.shapes.add_picture(str(parse_img), Inches(7.3), Inches(1.4),
                             width=Inches(5.6))
    add_footer(s, 5, total)

    # ------------------- 6. Mapping & Q-Learning -------------------
    s = prs.slides.add_slide(BLANK)
    add_title_block(s, "Mapping and Q-Learning")
    add_bullets(s, [
        "Per-maze state: trit edge maps (-1 unknown / 0 open / 1 wall).",
        "4 per-phase fire sets, confusion set, teleport source→dest mapping.",
        "Q[row, col, phase] — tabular risk, decayed between episodes.",
        "TD(0) update on death: propagates to Manhattan neighbors with geometric decay.",
        "What transfers between mazes: the RiskConfig hyperparameters (learning rate, decay, penalty scale).",
        "What resets between mazes: per-cell Q-values and the edge map.",
    ], left=0.6, top=1.4, width=12.1, height=5.5, font_size=19)
    add_footer(s, 6, total)

    # ------------------- 7. Weighted A* -------------------
    s = prs.slides.add_slide(BLANK)
    add_title_block(s, "Planner — Weighted A*")
    add_bullets(s, [
        "Edge cost = 1  +  fire penalty (at arrival phase)  +  confusion cost  +  w·Q-value  +  visit penalty",
        "Strict pass: block every cell that has ever been a fire.",
        "Relaxed pass: block only fires at the current arrival phase.",
        "Frontier fallback: if goal unreachable, head toward nearest frontier cell.",
        "Batching: append actions up to 5 while each next edge is known-open and landing is safe.",
    ], left=0.6, top=1.4, width=12.1, height=5.5, font_size=20)
    add_footer(s, 7, total)

    # ------------------- 8. Working Demo: alpha (NEW) ------------
    s = prs.slides.add_slide(BLANK)
    add_title_block(s, "Working Demo — maze-alpha:  Learning → Optimal")

    # Left caption ABOVE image
    cap = s.shapes.add_textbox(Inches(0.4), Inches(1.05),
                               Inches(6.2), Inches(0.9))
    tf = cap.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Episode 1 — RL Exploration"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "4796 turns • 16 deaths • agent learns the map from scratch"
    p.font.size = Pt(13); p.alignment = PP_ALIGN.CENTER

    # Left image (below its caption)
    gif_a1 = FIG / "alpha_demo_exploration.gif"
    fallback_a1 = FIG / "alpha_demo_exploration_final.png"
    traj_a = FIG / "trajectory_alpha.png"
    img_a1 = (gif_a1 if gif_a1.exists()
              else fallback_a1 if fallback_a1.exists()
              else traj_a if traj_a.exists() else None)
    if img_a1 is not None:
        s.shapes.add_picture(str(img_a1), Inches(0.7), Inches(2.1),
                             width=Inches(5.6))

    # Right caption ABOVE image
    cap = s.shapes.add_textbox(Inches(6.7), Inches(1.05),
                               Inches(6.2), Inches(0.9))
    tf = cap.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Episode 2+ — BFS-Optimal (Converged)"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "87 turns • 0 deaths • clean BFS path on the learned map"
    p.font.size = Pt(13); p.alignment = PP_ALIGN.CENTER

    # Right image (below its caption)
    gif_a2 = FIG / "alpha_demo_converged.gif"
    fallback_a2 = FIG / "alpha_demo_converged_final.png"
    img_a2 = (gif_a2 if gif_a2.exists()
              else fallback_a2 if fallback_a2.exists()
              else traj_a if traj_a.exists() else None)
    if img_a2 is not None:
        s.shapes.add_picture(str(img_a2), Inches(7.0), Inches(2.1),
                             width=Inches(5.6))

    add_footer(s, 8, total)

    # ------------------- 9. Working Demo: beta (existing) --------
    s = prs.slides.add_slide(BLANK)
    add_title_block(s, "Working Demo — maze-beta:  Learning → Optimal")

    # Left caption ABOVE image
    cap = s.shapes.add_textbox(Inches(0.4), Inches(1.05),
                               Inches(6.2), Inches(0.9))
    tf = cap.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Episode 1 — RL Exploration"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "5204 turns • 49 deaths • agent learns the map from scratch"
    p.font.size = Pt(13); p.alignment = PP_ALIGN.CENTER

    # Left image (below its caption)
    gif1 = FIG / "beta_demo_exploration.gif"
    fallback1 = FIG / "beta_demo_exploration_final.png"
    img1 = gif1 if gif1.exists() else fallback1
    if img1.exists():
        s.shapes.add_picture(str(img1), Inches(0.7), Inches(2.1),
                             width=Inches(5.6))

    # Right caption ABOVE image
    cap = s.shapes.add_textbox(Inches(6.7), Inches(1.05),
                               Inches(6.2), Inches(0.9))
    tf = cap.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Episode 3+ — BFS-Optimal (Converged)"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "~80 turns • 0 deaths • clean path via (2,32)→(1,32)→(0,32)→(0,31)"
    p.font.size = Pt(13); p.alignment = PP_ALIGN.CENTER

    # Right image (below its caption)
    gif2 = FIG / "beta_demo_converged.gif"
    fallback2 = FIG / "beta_demo_converged_final.png"
    traj2 = FIG / "trajectory_beta.png"
    if gif2.exists():
        img2 = gif2
    elif fallback2.exists():
        img2 = fallback2
    elif traj2.exists():
        img2 = traj2
    else:
        img2 = None
    if img2 is not None:
        s.shapes.add_picture(str(img2), Inches(7.0), Inches(2.1),
                             width=Inches(5.6))

    add_footer(s, 9, total)

    # ------------------- 9. Results -------------------
    s = prs.slides.add_slide(BLANK)
    add_title_block(s, "Results — alpha + beta (5 eval episodes each)")

    # Table
    rows = [
        ("Metric", "maze-alpha", "maze-beta"),
        ("Success rate", "100%", "100%"),
        ("Average turns", "1015.4", "1178.4"),
        ("Average path length", "2695.8", "3680.2"),
        ("Death rate", "0.0030", "0.0090"),
        ("Exploration efficiency", "0.83", "0.68"),
        ("Map completeness", "0.33", "0.30"),
        ("Replanning (≈ ms/turn)", "3", "50"),
        ("Learning efficiency (eps to 1st)", "1", "1"),
        ("Eps 2–5 converged turns", "87", "80"),
    ]
    table = s.shapes.add_table(len(rows), 3,
                                Inches(0.8), Inches(1.3),
                                Inches(11.7), Inches(5.3)).table
    for c, w in enumerate([Inches(5.3), Inches(3.2), Inches(3.2)]):
        table.columns[c].width = w
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(16 if r == 0 else 15)
            p.font.bold = (r == 0) or (c == 0)
    note = s.shapes.add_textbox(Inches(0.8), Inches(6.7),
                                Inches(11.7), Inches(0.5))
    p = note.text_frame.paragraphs[0]
    p.text = ("Both mazes meet the spec's stretch goals (>95% success, "
              "<500 turns on converged episodes, <1% deaths). "
              "No training on beta — spec requirement.")
    p.font.size = Pt(13)
    p.font.italic = True
    add_footer(s, 10, total)

    # ------------------- 10. Bonus Opportunities -------------------
    s = prs.slides.add_slide(BLANK)
    add_title_block(s, "Bonus Opportunities (spec §8.2, cap: 10 pts)")

    # Bullets on the left, dashboard image on the right.
    tbox = s.shapes.add_textbox(Inches(0.4), Inches(1.3),
                                Inches(5.3), Inches(5.7))
    tf = tbox.text_frame
    tf.word_wrap = True
    for i, (hdr, body) in enumerate([
        ("Maze-gamma attempt (5 pts)",
         "5-ep evaluation run, discovered 1 of 2 teleports, reached "
         "Manhattan 25 from goal; not solved in budget."),
        ("Novel visualization (up to 5 pts)",
         "src/solution_dashboard.py — unified 6-panel dashboard of "
         "maps, trajectories, convergence, metrics (shown →)."),
        ("Open-source toolkit release (up to 5 pts)",
         "MIT license, complete README with install + public API "
         "example, spec-agnostic code."),
    ]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(10)
        p.text = "• " + hdr
        p.font.size = Pt(18); p.font.bold = True
        p.font.color.rgb = ACCENT
        p2 = tf.add_paragraph()
        p2.text = "  " + body
        p2.font.size = Pt(14)
        p2.space_after = Pt(4)

    dashboard = FIG / "solution_dashboard.png"
    if dashboard.exists():
        # Right half: x≈6.1 to 13.0, y≈1.1 to ~7.0 — 6.9" wide
        # Place the dashboard centered in that box.
        s.shapes.add_picture(str(dashboard), Inches(6.0), Inches(1.1),
                             width=Inches(7.0))
    add_footer(s, 11, total)

    # ------------------- 11. What made beta work -------------------
    s = prs.slides.add_slide(BLANK)
    add_title_block(s,
                    "What Turned Beta from ‘Unsolvable' to 100%")
    add_bullets(s, [
        "Off-by-one phase fix: environment checks pits at pre-increment phase; we had been using post-increment, so every 5th planned step walked into a fire.",
        "Goal-biased frontier targeting: enumerate known-open cells that still have an unknown edge, sort by Manhattan distance to goal, BFS to the closest one.",
        "A* cycle guard: negative-cost teleport bonuses could make came_from cycle; the reconstruction loop spun forever. Floored edge costs at 0.05 and added a seen-set guard.",
    ], left=0.6, top=1.5, width=12.1, height=5.5, font_size=19)
    add_footer(s, 12, total)

    # ------------------- 12. Conclusion -------------------
    s = prs.slides.add_slide(BLANK)
    add_title_block(s, "Conclusion & Lessons")
    add_bullets(s, [
        "Hybrid RL + phase-aware BFS + frontier planner reaches 100% success on both alpha and beta (5/5 each, no training on beta).",
        "Episode 1 is exploratory; episodes 2–5 run the BFS-optimal path in 80–87 turns with 0 deaths.",
        "Beta's walled-off goal would have defeated pure nearest-frontier exploration — goal-biased frontier targeting is what made the transfer work.",
        "AI usage fully disclosed in report.pdf (Claude Code, Opus 4.7, 1M-context).",
    ], left=0.6, top=1.4, width=12.1, height=4.2, font_size=20)
    closer = s.shapes.add_textbox(Inches(0.6), Inches(6.1),
                                  SLIDE_W - Inches(1.2), Inches(0.9))
    p = closer.text_frame.paragraphs[0]
    p.text = "Thank you  —  Questions?"
    p.font.size = Pt(36); p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER
    add_footer(s, 13, total)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
